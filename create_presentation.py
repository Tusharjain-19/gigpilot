import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # completely blank layout

    # Color Palette - Website Aesthetic
    BG_COLOR = RGBColor(11, 15, 23)        # #0B0F17 Obsidian Slate
    CARD_BG = RGBColor(17, 24, 39)         # #111827 Surface Card
    CARD_BORDER = RGBColor(39, 42, 49)     # #272A31 Border
    ACCENT_GREEN = RGBColor(21, 128, 61)   # #15803D Brand Emerald
    MINT_GREEN = RGBColor(52, 211, 153)    # #34D399 Bright Mint
    ACCENT_ORANGE = RGBColor(194, 65, 12)   # #C2410C Surge Orange
    TEXT_WHITE = RGBColor(249, 250, 251)    # #F9FAFB Crisp White
    TEXT_MUTED = RGBColor(156, 163, 175)   # #9CA3AF Muted Gray

    def add_blank_slide_with_bg():
        slide = prs.slides.add_slide(blank_layout)
        # Background rectangle
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return slide

    def add_header(slide, title_text, category_text="HACKATHON PITCH DECK"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = MINT_GREEN

        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # =========================================================================
    # SLIDE 1: Title & Team Slide
    # =========================================================================
    slide1 = add_blank_slide_with_bg()
    
    # Hero Banner Card
    add_card(slide1, 0.8, 1.0, 11.733, 5.5, bg_color=CARD_BG, border_color=ACCENT_GREEN)
    
    # Title Text Box
    t_box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(10.9), Inches(2.2))
    tf1 = t_box1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "GIGPILOT AI"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = MINT_GREEN

    p2 = tf1.add_paragraph()
    p2.text = "Proactive AI Copilot & Rights Advisor for Gig-Workers"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "Domain: GigShield | Problem Statement Solution"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(6)

    # Team Members Container
    add_card(slide1, 1.2, 4.0, 10.933, 2.0, bg_color=RGBColor(14, 20, 31), border_color=CARD_BORDER)
    
    t_team = slide1.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.6))
    tf_team = t_team.text_frame
    tf_team.word_wrap = True
    
    pt1 = tf_team.paragraphs[0]
    pt1.text = "TEAM MEMBERS"
    pt1.font.size = Pt(11)
    pt1.font.bold = True
    pt1.font.color.rgb = MINT_GREEN
    
    # 3 Team Members
    members = [
      ("Ayush Kumar", "AI & Backend Lead"),
      ("Syed Ather", "Full Stack Developer"),
      ("Tushar Jain", "System Architecture & UI/UX")
    ]
    
    pt2 = tf_team.add_paragraph()
    pt2.space_before = Pt(8)
    for name, role in members:
        run_name = pt2.add_run()
        run_name.text = f"•  {name} "
        run_name.font.size = Pt(14)
        run_name.font.bold = True
        run_name.font.color.rgb = TEXT_WHITE
        
        run_role = pt2.add_run()
        run_role.text = f"({role})    "
        run_role.font.size = Pt(12)
        run_role.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: Problem Statement & Market Pain Points
    # =========================================================================
    slide2 = add_blank_slide_with_bg()
    add_header(slide2, "The Problem: Blind Navigation & Gig Exploitation", "DOMAIN: GIGSHIELD")

    problems = [
      ("Hidden Fuel & Idling Losses", "Orders look lucrative (₹120) but lose money after traffic delays and rising fuel costs (₹110.93-₹111.68/L in Bangalore).", ACCENT_ORANGE),
      ("Oversaturated Hotspots", "Workers flock to crowded restaurant hubs without knowing real-time driver competition density.", ACCENT_ORANGE),
      ("Burnout & Fatigue Risks", "Continuous 4+ hour driving without rest breaks degrades safety scores and increases road accidents.", ACCENT_ORANGE),
      ("Lack of Rate Transparency", "Gig platforms underpay per-km rates without clear fare benchmark visibility or legal rights awareness.", ACCENT_ORANGE)
    ]

    for idx, (title, desc, color) in enumerate(problems):
        row = idx // 2
        col = idx % 2
        left = 0.8 + col * 5.95
        top = 1.6 + row * 2.6
        
        add_card(slide2, left, top, 5.75, 2.4)
        tb = slide2.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.3), Inches(5.15), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"0{idx+1}. {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(8)

    # =========================================================================
    # SLIDE 3: Our Solution — GigPilot AI
    # =========================================================================
    slide3 = add_blank_slide_with_bg()
    add_header(slide3, "Our Solution: GigPilot AI Proactive Copilot", "OUR SOLUTION")

    # Big Banner Summary
    add_card(slide3, 0.8, 1.6, 11.733, 1.4, bg_color=RGBColor(21, 128, 61), border_color=MINT_GREEN)
    tb_sum = slide3.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(1.0))
    tf_sum = tb_sum.text_frame
    tf_sum.word_wrap = True
    p_sum = tf_sum.paragraphs[0]
    p_sum.text = "GigPilot AI turns gig work from a blind hustle into an optimized, high-yield micro-business. It provides real-time traffic decisioning, instant OCR order profit calculation, continuous burnout monitoring, and multilingual rights protection."
    p_sum.font.size = Pt(14)
    p_sum.font.bold = True
    p_sum.font.color.rgb = TEXT_WHITE

    pillars = [
      ("360° Live Opportunity Radar", "Analyzes driver competition & monsoon surge to steer riders to high-demand hotspots (₹365/hr)."),
      ("AI Screenshot OCR Analyzer", "Parses order slip photos (Swiggy/Zomato/Uber) to compute net profit after fuel expenses."),
      ("Burnout & Fatigue Guardian", "Monitors continuous active shift hours (≥4.0h) to trigger mandatory safety recovery breaks."),
      ("Multilingual Rights Companion", "Supports 7 regional Indian languages with Groq Llama 3.3-70B AI for fare benchmark clarity.")
    ]

    for idx, (title, desc) in enumerate(pillars):
        left = 0.8 + idx * 2.98
        top = 3.3
        add_card(slide3, left, top, 2.8, 3.6)
        
        tb = slide3.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(2.4), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Pillar 0{idx+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = MINT_GREEN
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_before = Pt(4)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(8)

    # =========================================================================
    # SLIDE 4: Core Features Implemented
    # =========================================================================
    slide4 = add_blank_slide_with_bg()
    add_header(slide4, "Features Implemented: Live Interactive System", "PRODUCT FEATURES")

    features = [
      ("Vector Map Radar & GPS", "Leaflet vector map centered at Basavanagudi (BMSCE) with 360° revolving conic radar sweep and live GPS button.", MINT_GREEN),
      ("Order Slip OCR Vision", "Drag-and-drop receipt upload powered by Tesseract.js to detect ₹ payout, km distance, and net profit.", MINT_GREEN),
      ("TomTom Live Traffic Engine", "Integrates TomTom Flow API to penalize gridlocked routes (+38% fuel idling penalty on heavy traffic).", ACCENT_ORANGE),
      ("Direct Google Maps Routing", "1-click 'Route' button launching turn-by-turn navigation directly to high-yield order zones.", MINT_GREEN),
      ("GigDNA Reputation Index", "5D score evaluating Reliability, Safety, Efficiency, Income Stability, and Customer Happiness.", MINT_GREEN),
      ("7 Regional Indian Languages", "Seamless switcher supporting English, Hindi, Kannada, Bengali, Marathi, Telugu, and Tamil.", MINT_GREEN)
    ]

    for idx, (ftitle, fdesc, fcol) in enumerate(features):
        row = idx // 3
        col = idx % 3
        left = 0.8 + col * 3.95
        top = 1.6 + row * 2.7
        
        add_card(slide4, left, top, 3.8, 2.5)
        tb = slide4.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.25), Inches(3.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ftitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = fcol
        
        pd = tf.add_paragraph()
        pd.text = fdesc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_before = Pt(6)

    # =========================================================================
    # SLIDE 5: Tech Stack & System Architecture
    # =========================================================================
    slide5 = add_blank_slide_with_bg()
    add_header(slide5, "Tech Stack & System Architecture", "TECHNICAL SPECS")

    stacks = [
      ("Frontend Layer", "React 18, Vite, Tailwind CSS, Lucide Icons, Leaflet Vector Maps, Tesseract.js OCR"),
      ("Backend Services", "Node.js, Express Framework, Axios, SQLite3, Dual In-Memory Fast Data Store"),
      ("AI & Machine Learning", "Groq Llama 3.3-70B LLM API, Tesseract OCR Engine, TomTom Traffic Flow AI Engine"),
      ("Third-Party APIs", "TomTom Traffic & Routing API, OpenWeatherMap API, Google Maps Directions API")
    ]

    for idx, (stitle, sdesc) in enumerate(stacks):
        left = 0.8 + idx * 2.98
        top = 1.6
        add_card(slide5, left, top, 2.8, 5.2)
        
        tb = slide5.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(2.4), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"LAYER 0{idx+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = MINT_GREEN
        
        pt = tf.add_paragraph()
        pt.text = stitle
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pt.space_before = Pt(6)
        
        pd = tf.add_paragraph()
        pd.text = sdesc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_before = Pt(12)

    # =========================================================================
    # SLIDE 6: AI Component & Integration Deep-Dive
    # =========================================================================
    slide6 = add_blank_slide_with_bg()
    add_header(slide6, "AI Component: Genuine Integration & Innovation", "AI & ML SPECS")

    ai_items = [
      ("Groq Llama 3.3-70B LLM", "Powers the multilingual Copilot chatbot. Answers queries on fuel rates (₹110.93-₹111.68/L), shift spending, earnings summaries, and fare benchmarks."),
      ("Tesseract.js Client Vision", "Extracts payout amount (₹), distance (km), platform name, and delivery location from uploaded order screenshots with 94%+ OCR accuracy."),
      ("TomTom & Weather AI Engine", "Dynamically computes traffic idling delays and monsoon rain surge multipliers (+25% bonus payout) in real-time.")
    ]

    for idx, (atitle, adesc) in enumerate(ai_items):
        top = 1.6 + idx * 1.8
        add_card(slide6, 0.8, top, 11.733, 1.6)
        
        tb = slide6.shapes.add_textbox(Inches(1.1), Inches(top + 0.25), Inches(11.133), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"AI MODEL 0{idx+1}: {atitle}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = MINT_GREEN
        
        pd = tf.add_paragraph()
        pd.text = adesc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_WHITE
        pd.space_before = Pt(4)

    # =========================================================================
    # SLIDE 7: Business Model & Market Economics
    # =========================================================================
    slide7 = add_blank_slide_with_bg()
    add_header(slide7, "Business Model & Monetization Strategy", "MARKET & ECONOMICS")

    models = [
      ("Freemium Worker Subscriptions", "Basic radar & AI decisioning free forever. Premium tier (₹49/month) unlocks predictive heatmaps & unlimited instant OCR slip analysis.", MINT_GREEN),
      ("B2B Fleet & Rental Partnerships", "Partnering with EV rental companies (Yulu, Zypp) and gig platforms for driver retention, safety analytics, and battery swapping routing.", MINT_GREEN),
      ("GigDNA Financial Services", "Utilizing 5D GigDNA reputation scores to partner with fintech providers for micro-loans, instant daily payouts, and health insurance.", MINT_GREEN)
    ]

    for idx, (mtitle, mdesc, mcol) in enumerate(models):
        left = 0.8 + idx * 3.95
        top = 1.6
        add_card(slide7, left, top, 3.8, 5.2)
        
        tb = slide7.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.3), Inches(3.2), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"STREAM 0{idx+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = MINT_GREEN
        
        pt = tf.add_paragraph()
        pt.text = mtitle
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pt.space_before = Pt(6)
        
        pd = tf.add_paragraph()
        pd.text = mdesc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_before = Pt(12)

    # =========================================================================
    # SLIDE 8: Worker Financial Impact & ROI
    # =========================================================================
    slide8 = add_blank_slide_with_bg()
    add_header(slide8, "Worker Financial Impact: Proven Earnings Boost", "ROI & IMPACT")

    # Left Card: Traditional Hustle
    add_card(slide8, 0.8, 1.6, 5.75, 5.2, bg_color=CARD_BG, border_color=ACCENT_ORANGE)
    tb_l = slide8.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.15), Inches(4.6))
    tfl = tb_l.text_frame
    tfl.word_wrap = True
    
    pl = tfl.paragraphs[0]
    pl.text = "TRADITIONAL GIG WORK (WITHOUT COPILOT)"
    pl.font.size = Pt(14)
    pl.font.bold = True
    pl.font.color.rgb = ACCENT_ORANGE
    
    metrics_before = [
      ("Avg. Hourly Yield", "₹140 / hr"),
      ("Daily Fuel Waste", "₹180 / day"),
      ("Underpaid Orders Taken", "35% of orders"),
      ("Net Daily Shift Profit", "₹450 / shift")
    ]
    for label, val in metrics_before:
        p = tfl.add_paragraph()
        p.text = f"•  {label}: {val}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(10)

    # Right Card: With GigPilot AI
    add_card(slide8, 6.75, 1.6, 5.75, 5.2, bg_color=CARD_BG, border_color=MINT_GREEN)
    tb_r = slide8.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.15), Inches(4.6))
    tfr = tb_r.text_frame
    tfr.word_wrap = True
    
    pr = tfr.paragraphs[0]
    pr.text = "WITH GIGPILOT AI COPILOT (+73% PROFIT)"
    pr.font.size = Pt(14)
    pr.font.bold = True
    pr.font.color.rgb = MINT_GREEN
    
    metrics_after = [
      ("Avg. Hourly Yield", "₹245 / hr (+75% gain!)"),
      ("Daily Fuel Saved", "₹85 saved / day"),
      ("Underpaid Orders Rejected", "100% filtered via OCR"),
      ("Net Daily Shift Profit", "₹780 / shift (+73% boost!)")
    ]
    for label, val in metrics_after:
        p = tfr.add_paragraph()
        p.text = f"•  {label}: {val}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 9: Scalability & Future Roadmap
    # =========================================================================
    slide9 = add_blank_slide_with_bg()
    add_header(slide9, "Scalability & Future Product Roadmap", "FUTURE VISION")

    phases = [
      ("Phase 1: Hackathon MVP", "Live Opportunity Radar, Leaflet vector maps, Basavanagudi BMSCE hotspot engine, OCR profit analyzer, TomTom traffic API, 7 languages.", MINT_GREEN),
      ("Phase 2: Multi-App Aggregator", "Unified floating copilot overlay auto-accepting highest yield orders across Swiggy, Zomato, Uber, and Blinkit concurrently.", MINT_GREEN),
      ("Phase 3: EV & Fintech Scale", "Integration with EV battery swapping stations and leveraging 5D GigDNA credit scoring for instant rider micro-insurance.", MINT_GREEN)
    ]

    for idx, (ptitle, pdesc, pcol) in enumerate(phases):
        top = 1.6 + idx * 1.8
        add_card(slide9, 0.8, top, 11.733, 1.6)
        
        tb = slide9.shapes.add_textbox(Inches(1.1), Inches(top + 0.25), Inches(11.133), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ptitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = pcol
        
        pd = tf.add_paragraph()
        pd.text = pdesc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_WHITE
        pd.space_before = Pt(4)

    # =========================================================================
    # SLIDE 10: Conclusion & GitHub Demo Link
    # =========================================================================
    slide10 = add_blank_slide_with_bg()
    
    add_card(slide10, 0.8, 1.2, 11.733, 5.1, bg_color=CARD_BG, border_color=ACCENT_GREEN)
    
    tb10 = slide10.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.933), Inches(4.3))
    tf10 = tb10.text_frame
    tf10.word_wrap = True
    
    p1 = tf10.paragraphs[0]
    p1.text = "EMPWERING INDIA'S GIG WORKERS"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = MINT_GREEN
    
    p2 = tf10.add_paragraph()
    p2.text = "GigPilot AI — The Ultimate Proactive Copilot"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(6)
    
    p3 = tf10.add_paragraph()
    p3.text = "Thank you for reviewing our solution for the GigShield Problem Statement!"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(12)
    
    p4 = tf10.add_paragraph()
    p4.text = "• GitHub Repository: https://github.com/Tusharjain-19/gigpilot.git\n• Team: Ayush Kumar, Syed Ather, Tushar Jain"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = MINT_GREEN
    p4.space_before = Pt(16)

    # Save output file
    output_filename = "GigPilot_AI_Hackathon_Pitch.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'")

if __name__ == "__main__":
    build_presentation()
